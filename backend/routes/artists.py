from flask import Blueprint, request, jsonify, send_file, render_template, url_for, send_from_directory
from ..app import db
from backend.models.artist import Artist
from backend.models import Channel, ChannelStat, News, Activity

from datetime import datetime, date
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from io import BytesIO
from werkzeug.utils import secure_filename
import os
import logging

from markdown import markdown # Import markdown library
import requests # Import requests library
from backend.utils.wikipedia_utils import get_artist_info_from_wikipedia # New import
from backend.utils.gemini_utils import search_artist_ai # Added Gemini import
from backend.utils.auth import require_role # Added import
from backend.utils.wordcloud_generator import generate_wordcloud_for_artist, generate_wordcloud_from_text # Import wordcloud generator

UPLOAD_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', '..', 'frontend', 'public', 'images', 'artists', 'profile')
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif'}

bp = Blueprint('artists', __name__)
logger = logging.getLogger(__name__)

@bp.route('/search_and_generate_sql', methods=['GET'])
def search_and_generate_sql():
    try:
        artist_name = request.args.get('artist_name')
        if not artist_name:
            return jsonify({'error': 'Artist name is required.'}), 400

        # Check if artist already exists in DB
        try:
            existing_artist = Artist.query.filter_by(name=artist_name).first()
        except Exception as db_err:
            logger.error(f"Database error while searching for artist {artist_name}: {db_err}", exc_info=True)
            return jsonify({'error': f'Database error: {str(db_err)}'}), 500

        artist_exists_flag = existing_artist is not None
        
        # Try AI Search first
        logger.info(f"Searching Gemini AI for artist: {artist_name} (is_update={artist_exists_flag})")
        ai_result = search_artist_ai(artist_name, is_update=artist_exists_flag)
        
        # Determine data source for wordcloud
        wordcloud_text = ""
        artist_data = {}
        sources = []
        sql_query = ""
        python_script = None

        if ai_result:
            logger.info(f"Gemini AI result retrieved for {artist_name}")
            artist_data = ai_result['profile']
            sql_query = ai_result['sql_query']
            python_script = ai_result['python_script']
            sources = ai_result['sources']
            
            # Combine summary and source titles for wordcloud
            wordcloud_text = (artist_data.get('wiki_summary') or "") + " " + " ".join([s['title'] for s in sources])

        else:
            # Fallback to Wikipedia
            logger.info(f"Gemini AI failed, falling back to Wikipedia for artist: {artist_name}")
            artist_data = get_artist_info_from_wikipedia(artist_name)
            
            if not artist_data:
                logger.warning(f"No Wikipedia data found for artist: {artist_name}")
                return jsonify({'message': f'Could not find information for {artist_name} on Wikipedia or AI.'}), 404

            try:
                if existing_artist:
                    sql_query = generate_update_sql(existing_artist.id, artist_data)
                else:
                    sql_query = generate_insert_sql(artist_data)
            except Exception as sql_err:
                logger.error(f"Error generating SQL for {artist_name}: {sql_err}", exc_info=True)
                return jsonify({'error': f'SQL Generation error: {str(sql_err)}'}), 500
            
            sources = artist_data.get('sources', [])
            wordcloud_text = (artist_data.get('wiki_summary') or "") + " " + " ".join([s['title'] for s in sources])

        # Generate Wordcloud
        wordcloud_image = None
        top_keywords = []
        try:
            if wordcloud_text:
                wordcloud_image, top_keywords = generate_wordcloud_from_text(wordcloud_text)
        except Exception as wc_err:
            logger.error(f"Error generating wordcloud: {wc_err}")
            # Non-critical, continue without wordcloud

        return jsonify({
            'sql_query': sql_query,
            'python_script': python_script,
            'artist_data': artist_data,
            'artist_exists': artist_exists_flag,
            'sources': sources,
            'wordcloud_image': wordcloud_image,
            'top_keywords': top_keywords
        })

    except Exception as e:
        logger.error(f"Unexpected error in search_and_generate_sql: {e}", exc_info=True)
        return jsonify({'error': f'Internal Server Error: {str(e)}'}), 500

def generate_insert_sql(data):
    columns = []
    values = []
    
    # Default values for fields not easily extracted or provided by user
    data.setdefault('eng_name', 'NULL')
    data.setdefault('birth_date', 'NULL')
    data.setdefault('height_cm', 'NULL')
    data.setdefault('debut_date', 'NULL')
    data.setdefault('debut_title', 'NULL')
    data.setdefault('recent_activity_category', 'UNKNOWN')
    data.setdefault('recent_activity_name', 'UNKNOWN')
    data.setdefault('genre', 'NULL')
    data.setdefault('agency_id', 'NULL')
    data.setdefault('current_agency_name', 'NULL')
    data.setdefault('nationality', 'NULL')
    data.setdefault('is_korean', 1) # Default to Korean
    data.setdefault('gender', 'NA') # Default to Not Applicable
    data.setdefault('status', 'ACTIVE')
    data.setdefault('category_id', 'NULL')
    data.setdefault('platform', 'NULL')
    data.setdefault('social_media_url', 'NULL')
    data.setdefault('profile_photo', 'NULL')
    data.setdefault('guarantee_krw', 'NULL')
    data.setdefault('wiki_summary', 'NULL')

    for key, value in data.items():
        # Ensure column names match the schema (e.g., birth_date -> birth_date)
        if key in ['name', 'eng_name', 'birth_date', 'height_cm', 'debut_date', 'debut_title',
                   'recent_activity_category', 'recent_activity_name', 'genre', 'agency_id',
                   'current_agency_name', 'nationality', 'is_korean', 'gender', 'status',
                   'category_id', 'platform', 'social_media_url', 'profile_photo',
                   'guarantee_krw', 'wiki_summary']:
            
            columns.append(key)
            sql_value = "" # Initialize sql_value
            if isinstance(value, str):
                sql_value = "'" + value.replace("'", "''") + "'"
            elif type(value).__name__ == 'date':
                sql_value = f"'{value.strftime('%Y-%m-%d')}'"
            elif value is None:
                sql_value = "NULL"
            else:
                sql_value = str(value)
            values.append(sql_value)

    return f"INSERT INTO first_ent.Artists ({', '.join(columns)}) VALUES ({', '.join(values)});"

def generate_update_sql(artist_id, data):
    set_clauses = []
    
    # Default values for fields not easily extracted or provided by user
    # These will be used if the key is not present in data
    default_values = {
        'eng_name': 'NULL',
        'birth_date': 'NULL',
        'height_cm': 'NULL',
        'debut_date': 'NULL',
        'debut_title': 'NULL',
        'recent_activity_category': 'UNKNOWN',
        'recent_activity_name': 'UNKNOWN',
        'genre': 'NULL',
        'agency_id': 'NULL',
        'current_agency_name': 'NULL',
        'nationality': 'NULL',
        'is_korean': 1, # Default to Korean
        'gender': 'NA', # Default to Not Applicable
        'status': 'ACTIVE',
        'category_id': 'NULL',
        'platform': 'NULL',
        'social_media_url': 'NULL',
        'profile_photo': 'NULL',
        'guarantee_krw': 'NULL',
        'wiki_summary': 'NULL'
    }

    for key, value in data.items():
        if key in ['name', 'eng_name', 'birth_date', 'height_cm', 'debut_date', 'debut_title',
                   'recent_activity_category', 'recent_activity_name', 'genre', 'agency_id',
                   'current_agency_name', 'nationality', 'is_korean', 'gender', 'status',
                   'category_id', 'platform', 'social_media_url', 'profile_photo',
                   'guarantee_krw', 'wiki_summary']:
            
            sql_value = ""
            if isinstance(value, str):
                sql_value = "'" + value.replace("'", "''") + "'"
            elif type(value).__name__ == 'date':
                sql_value = f"'{value.strftime('%Y-%m-%d')}'"
            elif value is None:
                sql_value = "NULL"
            else:
                sql_value = str(value)
            set_clauses.append(f"{key}={sql_value}")

    # Add default values for fields not returned by Wikipedia if they are not already in set_clauses
    for key, default_value in default_values.items():
        if key not in data and key in ['name', 'eng_name', 'birth_date', 'height_cm', 'debut_date', 'debut_title',
                   'recent_activity_category', 'recent_activity_name', 'genre', 'agency_id',
                   'current_agency_name', 'nationality', 'is_korean', 'gender', 'status',
                   'category_id', 'platform', 'social_media_url', 'profile_photo',
                   'guarantee_krw', 'wiki_summary']:
             set_clauses.append(f"{key}={default_value}")


    return f"UPDATE first_ent.Artists SET {', '.join(set_clauses)} WHERE id={artist_id};"


@bp.route('/test-image')
def test_image():
    try:
        return send_file(os.path.join(UPLOAD_FOLDER, 'park_seo_joon_2025.png'), mimetype='image/png')
    except Exception as e:
        return str(e), 500

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@bp.route('/profile-photos/<filename>')
def serve_profile_photo(filename):
    return send_from_directory(UPLOAD_FOLDER, filename)

@bp.route('/', methods=['GET'])
def get_artists():
    """모든 아티스트 조회 (페이지네이션 포함)"""
    search_query = request.args.get('query')
    filter_status = request.args.get('status')
    filter_gender = request.args.get('gender')
    filter_nationality = request.args.get('nationality')
    filter_genre = request.args.get('genre')
    filter_category = request.args.get('category_id', type=int)
    
    page = request.args.get('page', 1, type=int)
    per_page = request.args.get('per_page', 10, type=int)

    artists_query = Artist.query

    if search_query:
        search_pattern = f'%{search_query}%'
        artists_query = artists_query.filter(
            db.or_(
                Artist.name.ilike(search_pattern),
                Artist.genre.ilike(search_pattern),
                Artist.nationality.ilike(search_pattern)
            )
        )
    
    if filter_status:
        artists_query = artists_query.filter(Artist.status == filter_status)
    if filter_gender:
        artists_query = artists_query.filter(Artist.gender == filter_gender)
    if filter_nationality:
        # Check if filter_nationality is 'KOREAN' or 'FOREIGN' and map to is_korean
        if filter_nationality == 'KOREAN':
            artists_query = artists_query.filter(Artist.is_korean == True)
        elif filter_nationality == 'FOREIGN':
            artists_query = artists_query.filter(Artist.is_korean == False)
    if filter_genre:
        artists_query = artists_query.filter(Artist.genre.ilike(f'%{filter_genre}%'))
    if filter_category:
        artists_query = artists_query.filter(Artist.category_id == filter_category)

    # Apply pagination
    pagination = artists_query.paginate(page=page, per_page=per_page, error_out=False)
    artists = pagination.items
    
    return jsonify({
        'artists': [artist.to_dict() for artist in artists],
        'total': pagination.total,
        'pages': pagination.pages,
        'current_page': pagination.page
    })

@bp.route('/recent', methods=['GET'])
def get_recent_artists():
    """최근 등록된 아티스트 조회"""
    limit = request.args.get('limit', 5, type=int)
    recent_artists = Artist.query.order_by(Artist.debut_date.desc()).limit(limit).all()
    return jsonify([artist.to_dict() for artist in recent_artists])

@bp.route('/<int:artist_id>', methods=['GET'])
def get_artist(artist_id):
    """특정 아티스트 조회"""
    artist = Artist.query.get_or_404(artist_id)
    return jsonify(artist.to_dict())

@bp.route('/dashboard/stats', methods=['GET'])
def get_dashboard_stats():
    """대시보드 통계 조회"""
    total_artists = Artist.query.count()
    # Assuming 'Channel' model exists and has a 'status' or similar field
    # For now, let's mock active channels and followers/views
    # In a real scenario, this would involve more complex queries and potentially a Channel model
    active_channels = 0 # Placeholder
    total_followers = "0M" # Placeholder
    monthly_views = "0M" # Placeholder

    return jsonify([
        {"title": "총 아티스트", "value": str(total_artists), "change": "", "changeType": "positive", "icon": "Users", "color": "from-blue-500 to-cyan-500"},
        {"title": "활성 채널", "value": str(active_channels), "change": "", "changeType": "positive", "icon": "Radio", "color": "from-purple-500 to-pink-500"},
        {"title": "총 팔로워", "value": total_followers, "change": "", "changeType": "positive", "icon": "TrendingUp", "color": "from-green-500 to-emerald-500"},
        {"title": "월간 조회수", "value": monthly_views, "change": "", "changeType": "positive", "icon": "Activity", "color": "from-orange-500 to-red-500"}
    ])

@bp.route('/dashboard/channel-performance', methods=['GET'])
def get_channel_performance():
    """대시보드 채널 성과 조회"""
    # Placeholder data for now. In a real scenario, this would query Channel and ChannelStat models.
    return jsonify([
        {"platform": "Instagram", "growth": "+15%", "posts": 24},
        {"platform": "YouTube", "growth": "+8%", "posts": 12},
        {"platform": "TikTok", "growth": "+23%", "posts": 18}
    ])

@bp.route('/', methods=['POST'])
def create_artist():
    """새 아티스트 생성"""
    # Ensure upload folder exists
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)

    profile_photo_path = None
    if 'profile_photos' in request.files:
        files = request.files.getlist('profile_photos')
        if files:
            # Assuming only one profile photo for now, or taking the first one
            file = files[0]
            if file and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                file_path = os.path.join(UPLOAD_FOLDER, filename)
                file.save(file_path)
                profile_photo_path = f"/images/artists/profile/{filename}"

    # Get other form data
    data = request.form

    artist = Artist(
        name=data['name'],
        birth_date=datetime.strptime(data['birth_date'], '%Y-%m-%d').date(),
        height_cm=int(data['height_cm']) if data.get('height_cm') else None,
        debut_date=datetime.strptime(data['debut_date'], '%Y-%m-%d').date() if data.get('debut_date') else None,
        genre=data.get('genre'),
        agency_id=int(data['agency_id']) if data.get('agency_id') else None,
        nationality=data.get('nationality'),
        is_korean=data.get('is_korean') == 'true',
        gender=data.get('gender'),
        status=data.get('status'),
        category_id=int(data['category_id']) if data.get('category_id') else None,
        platform=data.get('platform'),
        social_media_url=data.get('social_media_url'),
        profile_photo=profile_photo_path,
        guarantee_krw=int(data['guarantee_krw']) if data.get('guarantee_krw') else None
    )

    db.session.add(artist)
    db.session.commit()

    return jsonify(artist.to_dict()), 201

@bp.route('/<int:artist_id>', methods=['PUT'])
def update_artist(artist_id):
    """아티스트 정보 수정"""
    artist = Artist.query.get_or_404(artist_id)

    # Ensure upload folder exists
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)

    profile_photo_path = artist.profile_photo # Keep existing photo if not updated
    if 'profile_photos' in request.files:
        files = request.files.getlist('profile_photos')
        if files:
            file = files[0]
            if file and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                file_path = os.path.join(UPLOAD_FOLDER, filename)
                file.save(file_path)
                profile_photo_path = f"/images/artists/profile/{filename}"

    # Get other form data
    data = request.form
    
    artist.name = data.get('name', artist.name)
    if data.get('birth_date'):
        artist.birth_date = datetime.strptime(data['birth_date'], '%Y-%m-%d').date()
    artist.height_cm = int(data['height_cm']) if data.get('height_cm') else None
    if data.get('debut_date'):
        artist.debut_date = datetime.strptime(data['debut_date'], '%Y-%m-%d').date()
    artist.genre = data.get('genre', artist.genre)
    artist.agency_id = int(data['agency_id']) if data.get('agency_id') else None
    artist.nationality = data.get('nationality', artist.nationality)
    artist.is_korean = data.get('is_korean') == 'true'
    artist.gender = data.get('gender', artist.gender)
    artist.status = data.get('status', artist.status)
    artist.category_id = int(data['category_id']) if data.get('category_id') else None
    artist.platform = data.get('platform', artist.platform)
    artist.social_media_url = data.get('social_media_url', artist.social_media_url)
    artist.profile_photo = profile_photo_path
    artist.guarantee_krw = int(data['guarantee_krw']) if data.get('guarantee_krw') else artist.guarantee_krw
    
    db.session.commit()
    
    return jsonify(artist.to_dict())

@bp.route('/report', methods=['POST'], strict_slashes=False)
def generate_artist_report():
    """선택된 아티스트에 대한 보고서 생성"""
    try:
        data = request.get_json()
        artist_ids = data.get('artist_ids', [])
        report_format = data.get('report_format', 'pdf') # 'pdf' or 'pptx'

        if not artist_ids:
            return jsonify({'error': '아티스트 ID가 제공되지 않았습니다.'}), 400

        artists = Artist.query.filter(Artist.id.in_(artist_ids)).all()

        if not artists:
            return jsonify({'error': '선택된 아티스트를 찾을 수 없습니다.'}), 404

        if report_format == 'pdf':
            PUPPETEER_SERVICE_URL = "http://localhost:3001/generate-pdf"
            artist_reports_data = []
            
            for artist in artists:
                # --- Fetch related data for each artist ---
                artist_channels = Channel.query.filter_by(artist_id=artist.id).all()
                channel_stats_data = []
                for channel in artist_channels:
                    latest_stat = ChannelStat.query.filter_by(channel_id=channel.id).order_by(ChannelStat.stat_date.desc()).first()
                    if latest_stat:
                        channel_stats_data.append({
                            'platform': channel.platform,
                            'name': channel.channel_name,
                            'url': channel.channel_url,
                            'follower_count': latest_stat.follower_count,
                            'engagement_rate': float(latest_stat.engagement_rate) if latest_stat.engagement_rate else 0.0
                        })
                
                recent_news_articles = News.query.filter_by(artist_id=artist.id).order_by(News.published_at.desc()).limit(4).all()
                
                activities_cf = []
                activities_broadcast = []
                try:
                    all_activities = Activity.query.filter_by(artist_id=artist.id).order_by(Activity.start_time.desc()).all()
                    
                    # Enhanced filtering for Korean/English activity types
                    activities_cf = [a for a in all_activities if a.activity_type and any(kw in a.activity_type.upper() for kw in ['CF', 'AD', '광고', 'BRAND'])]
                    activities_broadcast = [a for a in all_activities if a.activity_type and any(kw in a.activity_type.upper() for kw in ['DRAMA', 'MOVIE', 'TV', '방송', '영화', '드라마'])]
                except Exception as act_err:
                    logger.warning(f"Could not fetch activities for artist {artist.name}: {act_err}")
                    db.session.rollback() # Important: rollback to clear the failed state
                    # Continue without activities if there's a schema issue

                profile_photo_filename = os.path.basename(artist.profile_photo) if artist.profile_photo else 'default.png'
                profile_photo_url = url_for('artists.serve_profile_photo', filename=profile_photo_filename, _external=True)

                # Generate word cloud for news analysis
                wordcloud_base64, top_keywords = generate_wordcloud_for_artist(artist.id)

                artist_reports_data.append({
                    'artist': artist,
                    'channel_stats': channel_stats_data,
                    'news': recent_news_articles,
                    'activities_cf': activities_cf,
                    'activities_broadcast': activities_broadcast,
                    'profile_photo_url': profile_photo_url,
                    'wordcloud_image': wordcloud_base64,
                    'top_keywords': top_keywords
                })

            artist_names_str = ", ".join([a.name for a in artists])
            
            rendered_html = render_template(
                'report_template.html',
                report_title="Artist Performance Analysis",
                artist_names_summary=artist_names_str,
                generation_date=datetime.now().strftime('%Y-%m-%d'),
                author_name="theProjectCompany STRATEGIC ANALYSIS",
                css_path=url_for('static', filename='css/report_styles.css', _external=True),
                artist_reports=artist_reports_data,
                footer_text="theProjectCompany MANAGEMENT - CONFIDENTIAL"
            )

            response = requests.post(
                PUPPETEER_SERVICE_URL,
                json={'htmlContent': rendered_html},
                timeout=120
            )
            response.raise_for_status()
            
            pdf_buffer = BytesIO(response.content)
            return send_file(
                pdf_buffer,
                download_name=f'artist_report_{datetime.now().strftime("%Y%m%d")}.pdf',
                mimetype='application/pdf'
            )

        elif report_format == 'pptx':
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation()

            for artist in artists:
                # Cover Slide
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = artist.name
                slide.placeholders[1].text = "Artist Commercial Proposal"

                if artist.profile_photo:
                    image_full_path = os.path.join(UPLOAD_FOLDER, os.path.basename(artist.profile_photo))
                    if os.path.exists(image_full_path):
                        try:
                            slide.shapes.add_picture(image_full_path, Inches(3), Inches(3), Inches(4), Inches(4))
                        except: pass

                # Details Slide
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = f"{artist.name} - Professional Profile"
                
                # Simple bullet points for PPTX
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(4))
                tf = txBox.text_frame
                tf.text = f"• Gender: {artist.gender}\n• Nationality: {artist.nationality}\n• Agency: {artist.current_agency_name}\n• Debut: {artist.debut_date}"

            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return send_file(buffer, download_name=f'artist_report_{datetime.now().strftime("%Y%m%d")}.pptx', mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

        return jsonify({'error': '지원되지 않는 보고서 형식입니다.'}), 400

    except Exception as e:
        logger.error(f"Error in generate_artist_report: {e}", exc_info=True)
        return jsonify({'error': f'보고서 생성 중 시스템 오류가 발생했습니다: {str(e)}'}), 500

@bp.route('/<int:artist_id>', methods=['DELETE'])
def delete_artist(artist_id):
    """아티스트 삭제"""
    artist = Artist.query.get_or_404(artist_id)
    # Explicitly merge the object into the current session before deletion
    artist = db.session.merge(artist)
    db.session.delete(artist)
    db.session.commit()
    
    return jsonify({'message': '아티스트가 삭제되었습니다.'}), 200

@bp.route('/upsert-from-wiki', methods=['POST'], strict_slashes=False)
@require_role('admin')
def upsert_artist_from_wiki():
    data = request.get_json()
    artist_data = data.get('artist_data')

    if not artist_data or not artist_data.get('name'):
        return jsonify({'error': 'Artist data or name is missing.'}), 400

    artist_name = artist_data['name']
    existing_artist = Artist.query.filter_by(name=artist_name).first()

    try:
        if existing_artist:
            # Update existing artist
            artist = existing_artist
            message = f"아티스트 '{artist_name}' 정보가 업데이트되었습니다."
        else:
            # Create new artist
            artist = Artist()
            message = f"새로운 아티스트 '{artist_name}' 정보가 추가되었습니다."
            db.session.add(artist)

        # Map and convert data types for all fields
        artist.name = artist_data.get('name', artist.name)
        artist.eng_name = artist_data.get('eng_name', artist.eng_name)
        
        # Handle birth_date
        if artist_data.get('birth_date') and artist_data['birth_date'] != 'NULL':
            artist.birth_date = datetime.strptime(artist_data['birth_date'], '%Y-%m-%d').date()
        else:
            artist.birth_date = None # Explicitly set to None for NULL

        artist.height_cm = int(artist_data['height_cm']) if artist_data.get('height_cm') and artist_data['height_cm'] != 'NULL' else None
        
        # Handle debut_date
        if artist_data.get('debut_date') and artist_data['debut_date'] != 'NULL':
            artist.debut_date = datetime.strptime(artist_data['debut_date'], '%Y-%m-%d').date()
        else:
            artist.debut_date = None # Explicitly set to None for NULL

        artist.debut_title = artist_data.get('debut_title', artist.debut_title)
        artist.recent_activity_category = artist_data.get('recent_activity_category', artist.recent_activity_category)
        artist.recent_activity_name = artist_data.get('recent_activity_name', artist.recent_activity_name)
        artist.genre = artist_data.get('genre', artist.genre)
        artist.agency_id = int(artist_data['agency_id']) if artist_data.get('agency_id') and artist_data['agency_id'] != 'NULL' else None
        artist.current_agency_name = artist_data.get('current_agency_name', artist.current_agency_name)
        artist.nationality = artist_data.get('nationality', artist.nationality)
        artist.is_korean = artist_data.get('is_korean', artist.is_korean) in [True, 'True', 1, '1'] # Handle boolean
        artist.gender = artist_data.get('gender', artist.gender)
        artist.status = artist_data.get('status', artist.status)
        artist.category_id = int(artist_data['category_id']) if artist_data.get('category_id') and artist_data['category_id'] != 'NULL' else None
        artist.platform = artist_data.get('platform', artist.platform)
        artist.social_media_url = artist_data.get('social_media_url', artist.social_media_url)
        artist.profile_photo = artist_data.get('profile_photo', artist.profile_photo)
        artist.guarantee_krw = int(artist_data['guarantee_krw']) if artist_data.get('guarantee_krw') and artist_data['guarantee_krw'] != 'NULL' else None
        artist.wiki_summary = artist_data.get('wiki_summary', artist.wiki_summary)

        db.session.commit()
        return jsonify({'message': message, 'artist': artist.to_dict()}), 200

    except Exception as e:
        db.session.rollback()
        logger.error(f"Error upserting artist {artist_name} from wiki: {e}", exc_info=True)
        return jsonify({'error': f'DB 적용 중 오류 발생: {e}'}), 500

from backend.services.news_crawler import NewsCrawler # Import NewsCrawler

@bp.route('/<int:artist_id>/crawl-news', methods=['POST'])
def trigger_news_crawl(artist_id):
    """특정 아티스트에 대한 뉴스 크롤링을 수동으로 트리거"""
    artist = Artist.query.get_or_404(artist_id)
    crawler = NewsCrawler()
    news_items = crawler.search_news_for_artist(artist)
    saved_count = crawler.save_news_to_db(news_items, artist)
    return jsonify({'message': f'{artist.name}에 대한 {saved_count}개 뉴스 크롤링 및 저장 완료', 'saved_count': saved_count})